"""Multi-format file content extraction (docs extra, lazy imports).

Converts xlsx/csv/docx/pptx/pdf/json/xml/txt into markdown text + metadata.
All heavy dependencies (pandas, openpyxl, python-docx, python-pptx, pypdf)
are imported lazily so the core runtime still boots without the ``docs``
extra installed — callers get :class:`ExtractorUnavailable` with an
actionable message instead of an ImportError at startup.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

# ext -> kind
_KIND_BY_EXT = {
    ".xlsx": "spreadsheet",
    ".xlsm": "spreadsheet",
    ".xls": "spreadsheet",
    ".csv": "table",
    ".tsv": "table",
    ".docx": "document",
    ".pptx": "presentation",
    ".pdf": "pdf",
    ".json": "data",
    ".xml": "data",
    ".txt": "text",
    ".md": "text",
    ".markdown": "text",
}

PREVIEWABLE_KINDS = {"spreadsheet", "table", "document", "presentation", "pdf"}


class ExtractorUnavailable(RuntimeError):
    """The format is supported but its optional dependency is missing."""


class UnsupportedFormat(ValueError):
    """No extractor exists for this extension."""


@dataclass
class Extracted:
    kind: str
    markdown: str
    metadata: dict = field(default_factory=dict)


def classify(path: str | Path) -> str:
    """Return the kind for a path's extension, or ``"unknown"``."""
    return _KIND_BY_EXT.get(Path(path).suffix.lower(), "unknown")


def is_supported(path: str | Path) -> bool:
    return classify(path) != "unknown"


def _require(module: str, pkg: str):
    try:
        return __import__(module)
    except ImportError as e:  # pragma: no cover - depends on environment
        raise ExtractorUnavailable(
            f"解析 {pkg} 需要 docs 依赖（缺失模块 {module}）。"
            "安装: uv sync --extra docs"
        ) from e


# --------------------------------------------------------------------------
# pandas helpers (shared by spreadsheet / table)
# --------------------------------------------------------------------------

def _read_spreadsheet(path: Path, sheet: str | None = None):
    """Return (pandas, DataFrame or dict-of-frames). sheet=None → all sheets."""
    pd = _require("pandas", "Excel/CSV")
    engines = ["calamine", "openpyxl"]
    last_err: Exception | None = None
    kwargs: dict = {"sheet_name": sheet if sheet is not None else None}
    for eng in engines:
        try:
            return pd, pd.read_excel(path, engine=eng, **kwargs)
        except ImportError as e:
            last_err = e
            continue
        except ValueError:
            # calamine cannot do some things (e.g. styles); fall back
            last_err = None
            continue
    if last_err is not None:
        raise ExtractorUnavailable(
            "解析 Excel 需要 python-calamine 或 openpyxl。安装: uv sync --extra docs"
        ) from last_err
    # default engine as last resort
    return pd, pd.read_excel(path, **kwargs)


def read_table(path: str | Path, sep: str | None = None) -> "object":
    """Read a CSV/TSV into a DataFrame, robustly.

    More tolerant than a bare ``pd.read_csv``:
    - ``encoding="utf-8-sig"`` strips a BOM (Excel-saved CSVs carry one, which
      would otherwise corrupt the first column name);
    - ``on_bad_lines="skip"`` skips rows whose field count doesn't match the
      header instead of raising ``ParserError`` (a stray extra comma in one row
      shouldn't kill the whole preview/analysis);
    - falls back to ``latin-1`` for files that aren't valid UTF-8.
    ``sep`` defaults from the extension (.tsv → tab, else comma).
    """
    pd = _require("pandas", "CSV")
    if sep is None:
        sep = "\t" if Path(path).suffix.lower() == ".tsv" else ","
    common = {"sep": sep, "on_bad_lines": "skip"}
    try:
        return pd.read_csv(path, encoding="utf-8-sig", **common)
    except UnicodeDecodeError:
        return pd.read_csv(path, encoding="latin-1", **common)


def _cell(v) -> str:
    """Stringify a cell value for markdown/JSON (NaN/NaT → '')."""
    try:
        import pandas as pd

        if v is None or (isinstance(v, float) and pd.isna(v)) or v is pd.NaT:
            return ""
    except Exception:
        if v is None:
            return ""
    if hasattr(v, "isoformat"):
        return v.isoformat()
    s = str(v)
    return s


def df_to_markdown(df, max_rows: int = 50) -> str:
    """Render a DataFrame as a GitHub-flavoured markdown table (no tabulate)."""
    cols = [str(c) for c in df.columns]
    out = ["| " + " | ".join(cols) + " |", "|" + "---|" * len(cols)]
    for _, row in df.head(max_rows).iterrows():
        cells = [_cell(v).replace("|", "\\|").replace("\n", " ") for v in row]
        out.append("| " + " | ".join(cells) + " |")
    if len(df) > max_rows:
        out.append(f"\n…（仅显示前 {max_rows} 行，共 {len(df)} 行）")
    return "\n".join(out)


def _sheet_meta(pd, frames: dict) -> list[dict]:
    meta = []
    for name, df in frames.items():
        meta.append(
            {
                "name": str(name),
                "rows": int(len(df)),
                "cols": int(len(df.columns)),
                "columns": [
                    {"name": str(c), "dtype": str(df[c].dtype)} for c in df.columns
                ],
            }
        )
    return meta


def _extract_spreadsheet(path: Path, max_rows: int) -> Extracted:
    pd, frames = _read_spreadsheet(path)
    if not isinstance(frames, dict):
        frames = {path.stem: frames}
    parts = []
    for name, df in frames.items():
        parts.append(f"## Sheet: {name}（{len(df)} 行 × {len(df.columns)} 列）\n")
        parts.append(df_to_markdown(df, max_rows=max_rows))
        parts.append("")
    return Extracted(
        kind="spreadsheet",
        markdown="\n".join(parts),
        metadata={"sheets": _sheet_meta(pd, frames)},
    )


def _extract_table(path: Path, max_rows: int) -> Extracted:
    pd = _require("pandas", "CSV")
    try:
        df = read_table(path)
    except Exception as e:
        raise UnsupportedFormat(f"无法解析 CSV/TSV: {e}") from e
    md = f"## {path.name}（{len(df)} 行 × {len(df.columns)} 列）\n\n" + df_to_markdown(
        df, max_rows=max_rows
    )
    return Extracted(
        kind="table",
        markdown=md,
        metadata={"sheets": _sheet_meta(pd, {path.stem: df})},
    )


# --------------------------------------------------------------------------
# Word / PowerPoint / PDF
# --------------------------------------------------------------------------

def _extract_docx(path: Path, max_chars: int) -> Extracted:
    _require("docx", "Word")
    from docx import Document

    doc = Document(str(path))
    out: list[str] = []
    for p in doc.paragraphs:
        style = (p.style.name or "").lower() if p.style else ""
        text = p.text.strip()
        if not text:
            continue
        if style.startswith("heading"):
            lvl = "".join(ch for ch in style if ch.isdigit()) or "2"
            out.append("#" * min(int(lvl) + 1, 6) + " " + text)
        elif style.startswith("title"):
            out.append("# " + text)
        elif "list bullet" in style:
            out.append("- " + text)
        elif "list number" in style:
            out.append("1. " + text)
        else:
            out.append(text)
        out.append("")
    for t in doc.tables:
        rows = [
            [(c.text or "").replace("|", "\\|").replace("\n", " ") for c in r.cells]
            for r in t.rows
        ]
        if not rows:
            continue
        out.append("| " + " | ".join(rows[0]) + " |")
        out.append("|" + "---|" * len(rows[0]))
        for r in rows[1:]:
            out.append("| " + " | ".join(r) + " |")
        out.append("")
    md = "\n".join(out)[:max_chars]
    props = doc.core_properties
    return Extracted(
        kind="document",
        markdown=md,
        metadata={
            "title": props.title or "",
            "author": props.author or "",
            "created": props.created.isoformat() if props.created else "",
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables),
        },
    )


def _extract_pptx(path: Path, max_chars: int) -> Extracted:
    _require("pptx", "PowerPoint")
    from pptx import Presentation

    prs = Presentation(str(path))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## Slide {i}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(
                    p.text for p in shape.text_frame.paragraphs if p.text.strip()
                ).strip()
                if text:
                    out.append(text)
            if getattr(shape, "has_table", False):
                rows = [
                    [(c.text or "").replace("|", "\\|").replace("\n", " ") for c in r.cells]
                    for r in shape.table.rows
                ]
                if rows:
                    out.append("| " + " | ".join(rows[0]) + " |")
                    out.append("|" + "---|" * len(rows[0]))
                    for r in rows[1:]:
                        out.append("| " + " | ".join(r) + " |")
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                out.append(f"> 备注: {notes}")
        out.append("")
    md = "\n".join(out)[:max_chars]
    props = prs.core_properties
    return Extracted(
        kind="presentation",
        markdown=md,
        metadata={
            "title": props.title or "",
            "author": props.author or "",
            "created": props.created.isoformat() if props.created else "",
            "slides": len(prs.slides),
        },
    )


def _extract_pdf(path: Path, max_chars: int) -> Extracted:
    _require("pypdf", "PDF")
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    out: list[str] = []
    for i, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        if text:
            out.append(f"## Page {i}\n\n{text}\n")
    meta = reader.metadata or {}
    return Extracted(
        kind="pdf",
        markdown="\n".join(out)[:max_chars],
        metadata={
            "title": (meta.title or "") if meta else "",
            "author": (meta.author or "") if meta else "",
            "pages": len(reader.pages),
        },
    )


def _extract_text(path: Path, max_chars: int, kind: str) -> Extracted:
    raw = path.read_bytes()
    text = raw.decode("utf-8", errors="replace")[:max_chars]
    return Extracted(kind=kind, markdown=text, metadata={"bytes": len(raw)})


# --------------------------------------------------------------------------
# public API
# --------------------------------------------------------------------------

def extract(path: str | Path, max_rows: int = 50, max_chars: int = 200_000) -> Extracted:
    """Extract a file into markdown + metadata. Raises on unknown/unsupported."""
    p = Path(path).expanduser()
    if not p.is_file():
        raise FileNotFoundError(f"文件不存在: {p}")
    kind = classify(p)
    if kind == "spreadsheet":
        return _extract_spreadsheet(p, max_rows=max_rows)
    if kind == "table":
        return _extract_table(p, max_rows=max_rows)
    if kind == "document":
        return _extract_docx(p, max_chars=max_chars)
    if kind == "presentation":
        return _extract_pptx(p, max_chars=max_chars)
    if kind == "pdf":
        return _extract_pdf(p, max_chars=max_chars)
    if kind in ("data", "text"):
        return _extract_text(p, max_chars=max_chars, kind=kind)
    raise UnsupportedFormat(
        f"不支持的文件格式: {p.suffix}。支持: {', '.join(sorted(_KIND_BY_EXT))}"
    )


def schema_summary(path: str | Path, sample_rows: int = 5) -> dict:
    """Compact schema for prompt injection (spreadsheets/tables only).

    Returns ``{"kind", "sheets": [{name, rows, cols, columns:[{name,dtype}],
    sample: [[...]]}]}`` — small enough to embed in the agent context.
    """
    p = Path(path).expanduser()
    kind = classify(p)
    if kind == "spreadsheet":
        _, frames = _read_spreadsheet(p)
        if not isinstance(frames, dict):
            frames = {p.stem: frames}
        sheets = _sheet_meta(__import__("pandas"), frames)
        for s in sheets:
            df = frames[s["name"]]
            s["sample"] = [
                [_cell(v) for v in row] for _, row in df.head(sample_rows).iterrows()
            ]
        return {"kind": kind, "sheets": sheets}
    if kind == "table":
        ex = _extract_table(p, max_rows=sample_rows)
        sheets = ex.metadata.get("sheets", [])
        if sheets:
            df = read_table(p)
            sheets[0]["sample"] = [
                [_cell(v) for v in row] for _, row in df.head(sample_rows).iterrows()
            ]
        return {"kind": kind, "sheets": sheets}
    raise UnsupportedFormat(f"schema_summary 仅支持表格类文件，当前: {p.suffix}")
