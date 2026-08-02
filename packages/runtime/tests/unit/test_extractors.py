"""Unit tests: multi-format extractors (xlsx/csv/docx/pptx/pdf).

Fixtures are generated programmatically (no binary files in the repo).
"""

from __future__ import annotations

from pathlib import Path

import pytest

from ginno_runtime.files import extractors as ex

pytestmark = pytest.mark.unit


# --------------------------------------------------------------------------- #
# fixture builders
# --------------------------------------------------------------------------- #

def make_xlsx(p: Path, rows: int = 3) -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "销售"
    ws.append(["日期", "产品", "金额"])
    for i in range(rows):
        ws.append([f"2026-01-{i + 1:02d}", f"产品{i}", 100.5 * (i + 1)])
    ws2 = wb.create_sheet("说明")
    ws2.append(["备注"])
    ws2.append(["这是说明页"])
    wb.save(p)
    return p


def make_csv(p: Path) -> Path:
    p.write_text("a,b\n1,2\n3,4\n", encoding="utf-8")
    return p


def make_docx(p: Path) -> Path:
    from docx import Document

    doc = Document()
    doc.core_properties.author = "张三"
    doc.core_properties.title = "测试合同"
    doc.add_heading("第一章", level=1)
    doc.add_paragraph("甲方与乙方签订本协议。")
    doc.add_paragraph("要点一", style="List Bullet")
    t = doc.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "项目"
    t.cell(0, 1).text = "金额"
    t.cell(1, 0).text = "A"
    t.cell(1, 1).text = "100"
    doc.save(p)
    return p


def make_pptx(p: Path) -> Path:
    from pptx import Presentation

    prs = Presentation()
    prs.core_properties.title = "产品演示"
    prs.core_properties.author = "李四"
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "欢迎"
    tb = slide.shapes.add_textbox(0, 0, 100, 100)
    tb.text_frame.text = "正文内容"
    slide.notes_slide.notes_text_frame.text = "演讲备注"
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.save(p)
    return p


def _minimal_pdf_bytes(text: str) -> bytes:
    """Hand-assemble a tiny valid PDF with one line of text (no deps)."""
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        None,  # content stream, filled below
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    content = f"BT /F1 24 Tf 72 720 Td ({text}) Tj ET".encode()
    objects[3] = b"<< /Length %d >>\nstream\n%s\nendstream" % (len(content), content)

    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, obj in enumerate(objects, 1):
        offsets.append(len(out))
        out += b"%d 0 obj\n%s\nendobj\n" % (i, obj)
    xref_pos = len(out)
    out += b"xref\n0 %d\n" % (len(objects) + 1)
    out += b"0000000000 65535 f \n"
    for off in offsets:
        out += b"%010d 00000 n \n" % off
    out += b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF\n" % (
        len(objects) + 1,
        xref_pos,
    )
    return bytes(out)


def make_pdf(p: Path) -> Path:
    p.write_bytes(_minimal_pdf_bytes("Hello Ginno"))
    return p


# --------------------------------------------------------------------------- #
# classify / dispatch
# --------------------------------------------------------------------------- #

def test_classify():
    assert ex.classify("a.xlsx") == "spreadsheet"
    assert ex.classify("a.CSV") == "table"
    assert ex.classify("a.docx") == "document"
    assert ex.classify("a.pptx") == "presentation"
    assert ex.classify("a.pdf") == "pdf"
    assert ex.classify("a.unknown") == "unknown"
    assert ex.is_supported("a.xls")
    assert not ex.is_supported("a.exe")


def test_unsupported_raises(tmp_path):
    f = tmp_path / "x.exe"
    f.write_bytes(b"MZ")
    with pytest.raises(ex.UnsupportedFormat):
        ex.extract(f)


def test_missing_file_raises():
    with pytest.raises(FileNotFoundError):
        ex.extract("/no/such/file.xlsx")


# --------------------------------------------------------------------------- #
# extraction
# --------------------------------------------------------------------------- #

def test_extract_xlsx(tmp_path):
    f = make_xlsx(tmp_path / "s.xlsx", rows=3)
    res = ex.extract(f)
    assert res.kind == "spreadsheet"
    assert "## Sheet: 销售" in res.markdown
    assert "| 日期 | 产品 | 金额 |" in res.markdown
    assert "产品0" in res.markdown
    assert "## Sheet: 说明" in res.markdown
    names = {s["name"] for s in res.metadata["sheets"]}
    assert names == {"销售", "说明"}
    sales = next(s for s in res.metadata["sheets"] if s["name"] == "销售")
    assert sales["rows"] == 3 and sales["cols"] == 3


def test_extract_csv(tmp_path):
    f = make_csv(tmp_path / "t.csv")
    res = ex.extract(f)
    assert res.kind == "table"
    assert "| a | b |" in res.markdown
    assert "| 1 | 2 |" in res.markdown
    assert res.metadata["sheets"][0]["rows"] == 2


def test_extract_docx(tmp_path):
    f = make_docx(tmp_path / "c.docx")
    res = ex.extract(f)
    assert res.kind == "document"
    assert "## 第一章" in res.markdown  # heading 1 → "##" (h1 reserved for title)
    assert "甲方与乙方签订本协议。" in res.markdown
    assert "- 要点一" in res.markdown
    assert "| 项目 | 金额 |" in res.markdown
    assert res.metadata["author"] == "张三"
    assert res.metadata["title"] == "测试合同"
    assert res.metadata["tables"] == 1


def test_extract_pptx(tmp_path):
    f = make_pptx(tmp_path / "d.pptx")
    res = ex.extract(f)
    assert res.kind == "presentation"
    assert "## Slide 1" in res.markdown
    assert "## Slide 2" in res.markdown
    assert "欢迎" in res.markdown
    assert "正文内容" in res.markdown
    assert "> 备注: 演讲备注" in res.markdown
    assert res.metadata["slides"] == 2
    assert res.metadata["title"] == "产品演示"


def test_extract_pdf(tmp_path):
    f = make_pdf(tmp_path / "x.pdf")
    res = ex.extract(f)
    assert res.kind == "pdf"
    assert "Hello Ginno" in res.markdown
    assert res.metadata["pages"] == 1


def test_extract_text(tmp_path):
    f = tmp_path / "n.txt"
    f.write_text("plain text", encoding="utf-8")
    res = ex.extract(f)
    assert res.kind == "text"
    assert res.markdown == "plain text"


# --------------------------------------------------------------------------- #
# schema summary + markdown rendering
# --------------------------------------------------------------------------- #

def test_schema_summary_xlsx(tmp_path):
    f = make_xlsx(tmp_path / "s.xlsx", rows=10)
    s = ex.schema_summary(f, sample_rows=2)
    assert s["kind"] == "spreadsheet"
    sales = next(x for x in s["sheets"] if x["name"] == "销售")
    assert sales["rows"] == 10 and sales["cols"] == 3
    cols = {c["name"]: c["dtype"] for c in sales["columns"]}
    assert cols["金额"] == "float64"
    assert len(sales["sample"]) == 2
    assert sales["sample"][0][1] == "产品0"


def test_schema_summary_csv(tmp_path):
    f = make_csv(tmp_path / "t.csv")
    s = ex.schema_summary(f)
    assert s["kind"] == "table"
    assert s["sheets"][0]["columns"][0]["name"] == "a"
    assert s["sheets"][0]["sample"][0] == ["1", "2"]


def test_schema_summary_rejects_doc(tmp_path):
    f = make_docx(tmp_path / "c.docx")
    with pytest.raises(ex.UnsupportedFormat):
        ex.schema_summary(f)


def test_df_to_markdown_handles_nan(tmp_path):
    pd = pytest.importorskip("pandas")
    df = pd.DataFrame({"a": [1, None], "b": ["x|y", "z\nw"]})
    md = ex.df_to_markdown(df)
    assert "| a | b |" in md
    assert "| 1.0 | x\\|y |" in md  # float col keeps .0; pipe escaped
    assert "|  | z w |" in md  # NaN → empty, newline → space
